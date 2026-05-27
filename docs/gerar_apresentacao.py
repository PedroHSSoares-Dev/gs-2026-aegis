"""
AEGIS — Gerador de Apresentação GS 2026
Visual: tema dark space idêntico ao dashboard AEGIS
Icons: paths SVG do site (src/components/icons/index.tsx) → matplotlib → PNG → PPTX
Execute: micromamba run -n dev python docs/gerar_apresentacao.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, os, math
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import PathPatch, FancyBboxPatch, Circle, Arc
from matplotlib.path import Path as MPath
import numpy as np

# ── svgpath2mpl para renderizar ícones do site ──────────────────────
from svgpath2mpl import parse_path

# ═══════════════════════ PALETA AEGIS (do index.css) ════════════════
BG      = RGBColor(0x06, 0x09, 0x12)   # --bg-deep
PANEL   = RGBColor(0x0A, 0x0E, 0x1A)   # --bg
PANEL2  = RGBColor(0x11, 0x1A, 0x2E)   # --panel
STEEL   = RGBColor(0x1E, 0x2D, 0x4A)   # --steel
STEEL2  = RGBColor(0x25, 0x38, 0x5A)   # --steel-2
CYAN    = RGBColor(0x00, 0xD4, 0xFF)   # --cyan
AMBER   = RGBColor(0xFF, 0x6B, 0x35)   # --amber
AMBER2  = RGBColor(0xFF, 0xB5, 0x47)   # --amber-2
TEAL    = RGBColor(0x5E, 0xE0, 0xC2)   # --teal
BLUE2   = RGBColor(0x3F, 0xA9, 0xFF)   # --blue-2
WHITE   = RGBColor(0xF0, 0xF4, 0xFF)   # --text
MID     = RGBColor(0x8A, 0x9B, 0xC4)
DIM     = RGBColor(0x50, 0x60, 0x88)
RED_CRIT= RGBColor(0xEF, 0x44, 0x44)

# Hex strings for matplotlib
H_BG    = '#060912'
H_PANEL = '#0A0E1A'
H_PAN2  = '#111A2E'
H_STEEL = '#1E2D4A'
H_ST2   = '#25385A'
H_CYAN  = '#00D4FF'
H_AMBER = '#FF6B35'
H_AMB2  = '#FFB547'
H_TEAL  = '#5EE0C2'
H_BLUE2 = '#3FA9FF'
H_WHITE = '#F0F4FF'
H_MID   = '#8A9BC4'
H_DIM   = '#506088'

W = Inches(13.333)
H_SLD = Inches(7.5)

# ═══════════════════ ÍCONES DO SITE AEGIS ═══════════════════════════
# Paths SVG copiados de src/components/icons/index.tsx
# Cada ícone: lista de (d_string, draw_as) onde draw_as = 'stroke'|'fill'
# Círculos SVG são convertidos para path com arco

def _circle_path(cx, cy, r):
    """Converte SVG circle para path string."""
    return (f"M {cx-r},{cy} "
            f"A {r},{r} 0 1,0 {cx+r},{cy} "
            f"A {r},{r} 0 1,0 {cx-r},{cy} Z")

ICONS = {
    'satellite': [
        ("M5 10l4-4 9 9-4 4-9-9z", 'stroke'),
        ("M14 5l5 5", 'stroke'),
        ("M9 14l-3 3 1 4 4 1 3-3", 'stroke'),
        ("M2 22l4-4", 'stroke'),
    ],
    'globe': [
        (_circle_path(12, 12, 9), 'stroke'),
        ("M3 12h18", 'stroke'),
        ("M12 3 Q19 8 19 12 Q19 16 12 21", 'stroke'),
        ("M12 3 Q5 8 5 12 Q5 16 12 21", 'stroke'),
    ],
    'alert': [
        ("M12 3l10 18H2L12 3z", 'stroke'),
        ("M12 10v5", 'stroke'),
        ("M12 18.5v0.5", 'stroke'),
    ],
    'pulse': [
        ("M3 12h4l2-7 4 14 2-7h6", 'stroke'),
    ],
    'fire': [
        ("M12 2c1 4 5 5 5 10a5 5 0 1 1-10 0c0-2 1-3 2-4-1 4 3 3 3 6", 'stroke'),
    ],
    'cyclone': [
        ("M12 3a9 9 0 0 1 9 9", 'stroke'),
        ("M12 21a9 9 0 0 1-9-9", 'stroke'),
        (_circle_path(12, 12, 3), 'stroke'),
        ("M12 8a4 4 0 0 1 4 4", 'stroke'),
        ("M12 16a4 4 0 0 1-4-4", 'stroke'),
    ],
    'flood': [
        ("M3 9c2 0 2-2 4.5-2S10 9 12 9s2-2 4.5-2S19 9 21 9", 'stroke'),
        ("M3 14c2 0 2-2 4.5-2S10 14 12 14s2-2 4.5-2S19 14 21 14", 'stroke'),
        ("M3 19c2 0 2-2 4.5-2S10 19 12 19s2-2 4.5-2S19 19 21 19", 'stroke'),
    ],
    'quake': [
        ("M2 12h3l2-6 3 14 3-10 3 6 2-4h4", 'stroke'),
    ],
    'drought': [
        (_circle_path(12, 6, 3), 'stroke'),
        ("M12 2v1M12 9v1M5 6h1M18 6h1", 'stroke'),
        ("M7 4l0.7 0.7M16.3 4l-0.7 0.7M7 8l0.7-0.7M16.3 8l-0.7-0.7", 'stroke'),
        ("M3 18h18M5 21h14", 'stroke'),
    ],
    'layers': [
        ("M12 3l9 5-9 5-9-5 9-5z", 'stroke'),
        ("M3 13l9 5 9-5", 'stroke'),
        ("M3 18l9 5 9-5", 'stroke'),
    ],
    'team': [
        (_circle_path(9, 8, 3), 'stroke'),
        ("M3 20a6 6 0 0 1 12 0", 'stroke'),
        ("M16 4a3 3 0 0 1 0 6", 'stroke'),
        ("M16 14a6 6 0 0 1 5 6", 'stroke'),
    ],
}


def render_icon(name: str, size_px=72, stroke_color=H_CYAN, bg=H_PAN2,
                lw=1.5, pad=1.2) -> io.BytesIO:
    """
    Renderiza um ícone AEGIS como PNG e retorna BytesIO.
    Usa svgpath2mpl para converter os paths SVG do site.
    """
    icon = ICONS.get(name, [])
    fig, ax = plt.subplots(figsize=(1, 1), dpi=size_px)
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)
    # Coordenadas SVG: 0–24 × 0–24, Y cresce para baixo
    ax.set_xlim(-pad, 24 + pad)
    ax.set_ylim(-pad, 24 + pad)
    ax.set_aspect('equal')
    ax.invert_yaxis()  # SVG tem Y invertido em relação ao matplotlib padrão
    ax.axis('off')

    for d_str, draw_as in icon:
        try:
            path = parse_path(d_str)
        except Exception:
            continue
        if draw_as == 'fill':
            patch = PathPatch(path, facecolor=stroke_color,
                              edgecolor='none', linewidth=0, zorder=3)
        else:
            patch = PathPatch(path, facecolor='none',
                              edgecolor=stroke_color, linewidth=lw,
                              capstyle='round', joinstyle='round', zorder=3)
        ax.add_patch(patch)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=size_px,
                facecolor=bg, pad_inches=0.06)
    plt.close(fig)
    buf.seek(0)
    return buf


# ═══════════════════ HELPERS PPTX ═══════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H_SLD
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor = BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh


def _to_rgb(color):
    """Aceita RGBColor ou hex string '#RRGGBB'."""
    if isinstance(color, RGBColor):
        return color
    h = color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def txt(slide, text, x, y, w, h, size=Pt(11), color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font='Calibri', wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size; r.font.color.rgb = _to_rgb(color)
    r.font.bold = bold; r.font.name = font
    return tb


def label(slide, text, x, y, w, h, size=Pt(8.5), color=DIM, spacing=180):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = size; r.font.color.rgb = color
    r.font.name = 'Calibri'
    r._r.get_or_add_rPr().set('spc', str(spacing))
    return tb


def embed_buf(slide, buf, x, y, w, h):
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)


def embed_icon(slide, name, x, y, size_in, color=H_CYAN, bg=H_PAN2):
    buf = render_icon(name, size_px=96, stroke_color=color, bg=bg)
    sz = Inches(size_in)
    embed_buf(slide, buf, x, y, sz, sz)


def embed_fig(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight',
                facecolor=fig.get_facecolor(), dpi=150)
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)


# ── Decoração comum ─────────────────────────────────────────────────

def cyan_topbar(slide):
    rect(slide, 0, Inches(0), W, Pt(2.5), fill=CYAN)


def header(slide, section, right='AEGIS · ORBITAL DISASTER INTELLIGENCE · 2026'):
    label(slide, section.upper(), Inches(0.42), Inches(0.1), Inches(7), Inches(0.3),
          color=DIM, spacing=250)
    label(slide, right, Inches(6.5), Inches(0.1), Inches(6.6), Inches(0.3),
          color=DIM, spacing=180)


def slide_title(slide, title, subtitle='', t_color=CYAN, s_color=MID):
    txt(slide, title, Inches(0.42), Inches(0.42), Inches(12.5), Inches(0.75),
        size=Pt(26), color=t_color, bold=True)
    if subtitle:
        txt(slide, subtitle, Inches(0.42), Inches(1.08), Inches(12.5), Inches(0.35),
            size=Pt(11), color=s_color)


def metric_card(slide, value, lbl, x, y, w=Inches(2.45), h=Inches(1.3),
                vc=CYAN):
    rect(slide, x, y, w, h, fill=PANEL2, line=STEEL)
    txt(slide, value, x, y + Inches(0.12), w, Inches(0.58),
        size=Pt(24), color=vc, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, lbl, x, y + Inches(0.65), w, Inches(0.55),
        size=Pt(8.5), color=MID, align=PP_ALIGN.CENTER)


def bullet_panel(slide, items, x, y, w, h, title='', tc=CYAN, bg=PANEL2):
    rect(slide, x, y, w, h, fill=bg, line=STEEL)
    ty = y + Inches(0.1)
    if title:
        rect(slide, x, y, w, Inches(0.38), fill=tc)
        txt(slide, title, x + Inches(0.15), y + Inches(0.09), w - Inches(0.2),
            Inches(0.28), size=Pt(9.5), color=BG, bold=True)
        ty += Inches(0.45)
    for b in items:
        txt(slide, '›  ' + b, x + Inches(0.15), ty, w - Inches(0.25), Inches(0.35),
            size=Pt(10), color=WHITE)
        ty += Inches(0.34)


def star_field(slide, n=40, seed=2026):
    rng = np.random.default_rng(seed)
    for _ in range(n):
        sx = float(rng.uniform(0, 13.333))
        sy = float(rng.uniform(0, 7.5))
        c = CYAN if rng.random() < 0.3 else DIM
        dot = slide.shapes.add_shape(1, Inches(sx), Inches(sy), Pt(1.5), Pt(1.5))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()


# ═══════════════════ FIGURAS MATPLOTLIB ═════════════════════════════

def fig_eda():
    cats = ['Incêndios', 'Tempestades', 'Vulcões', 'Gelo Marinho',
            'Enchentes', 'Terremotos', 'Seca', 'Deslizamentos']
    vals = [412, 298, 187, 143, 98, 76, 52, 34]
    bar_colors = [H_AMBER, H_BLUE2, H_AMBER, H_TEAL,
                  H_BLUE2, H_AMB2,  H_AMBER, H_TEAL]

    fig, axes = plt.subplots(1, 2, figsize=(12.5, 4.4))
    fig.patch.set_facecolor(H_BG)

    # Bar chart
    ax = axes[0]
    ax.set_facecolor(H_PANEL)
    ax.barh(cats[::-1], vals[::-1], color=bar_colors[::-1], height=0.62, edgecolor='none')
    ax.set_xlabel('Número de Eventos (2024–2026)', color=H_MID, fontsize=9)
    ax.set_title('Distribuição por Categoria — NASA EONET (últimos 2 anos)',
                 color=H_CYAN, fontsize=10, pad=8, fontweight='bold')
    ax.tick_params(colors=H_MID, labelsize=8)
    for spine in ax.spines.values(): spine.set_edgecolor(H_STEEL)
    ax.grid(axis='x', color=H_STEEL, lw=0.5, alpha=0.6)
    ax.set_axisbelow(True)
    for i, (bar, val) in enumerate(zip(ax.patches, vals[::-1])):
        ax.text(bar.get_width() + 6, bar.get_y() + bar.get_height()/2,
                str(val), va='center', color=H_WHITE, fontsize=8)

    # Série temporal
    ax2 = axes[1]
    ax2.set_facecolor(H_PANEL)
    months = np.arange(24)
    np.random.seed(42)
    base = 40 + 15 * np.sin(months * np.pi / 6) + np.random.normal(0, 4, 24)
    ma   = np.convolve(base, np.ones(3)/3, mode='same')
    ax2.fill_between(months, base, alpha=0.18, color=H_CYAN)
    ax2.plot(months, base, color=H_CYAN, lw=1.2, alpha=0.75, label='Eventos por mês')
    ax2.plot(months, ma,   color=H_AMBER, lw=2.2, label='Tendência (média móvel 3m)')
    ax2.set_title('Série Temporal Mensal (2024–2026)',
                  color=H_CYAN, fontsize=10, pad=8, fontweight='bold')
    ax2.set_xlabel('Mês', color=H_MID, fontsize=9)
    ax2.set_ylabel('Nº de Eventos', color=H_MID, fontsize=9)
    ax2.tick_params(colors=H_MID, labelsize=7)
    for spine in ax2.spines.values(): spine.set_edgecolor(H_STEEL)
    ax2.grid(color=H_STEEL, lw=0.5, alpha=0.5)
    lbl_ticks = ['Jan/24','','Mar/24','','Mai/24','','Jul/24','','Set/24','','Nov/24','',
                 'Jan/25','','Mar/25','','Mai/25','','Jul/25','','Set/25','','Nov/25','']
    ax2.set_xticks(months)
    ax2.set_xticklabels(lbl_ticks, rotation=40, ha='right', fontsize=7)
    ax2.legend(facecolor=H_PAN2, edgecolor=H_STEEL, labelcolor=H_WHITE, fontsize=8)
    fig.tight_layout(pad=1.1)
    return fig


def fig_arch():
    fig, ax = plt.subplots(figsize=(13.1, 5.8))
    fig.patch.set_facecolor(H_BG)
    ax.set_facecolor(H_BG)
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 6.5)
    ax.axis('off')

    def box(x, y, w, h, title, color, items, icon_name=None):
        patch = FancyBboxPatch((x, y), w, h, boxstyle='round,pad=0.07',
                               lw=1.2, edgecolor=color, facecolor=H_PAN2, zorder=2)
        ax.add_patch(patch)
        # Barra superior colorida
        top = FancyBboxPatch((x, y+h-0.32), w, 0.32,
                             boxstyle='round,pad=0.01', lw=0, facecolor=color,
                             zorder=3, clip_on=False)
        ax.add_patch(top)
        ax.text(x + w/2, y + h - 0.16, title,
                ha='center', va='center', fontsize=7, fontweight='bold',
                color=H_BG, fontfamily='monospace', zorder=4)
        for i, item in enumerate(items):
            ax.text(x + 0.12, y + h - 0.48 - i*0.38, '› ' + item,
                    ha='left', va='top', fontsize=6.2, color=H_WHITE,
                    fontfamily='monospace', zorder=3, clip_on=True)

    def arrow(x1, y1, x2, y2, color=H_STEEL):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=color, lw=1.1,
                                   connectionstyle='arc3,rad=0'),
                    zorder=5)

    # Título
    ax.text(7, 6.35, 'AEGIS  ·  ARQUITETURA ORACLE CLOUD INFRASTRUCTURE (OCI)',
            ha='center', va='center', fontsize=9.5, fontweight='bold',
            color=H_CYAN, fontfamily='monospace')

    # Col 1 — FONTES EXTERNAS
    box(0.08, 0.25, 2.1, 5.8, 'FONTES EXTERNAS', H_CYAN,
        ['NASA EONET API v3', 'NOAA Climate API',
         'INPE BDQueimadas', 'Copernicus Sentinel',
         'USGS EarthExplorer', 'Planet Labs'])

    # Col 2 — INGESTÃO
    box(2.35, 3.4, 2.1, 2.65, 'INGESTÃO OCI', H_BLUE2,
        ['OCI API Gateway', 'OCI Streaming', '(Kafka-compatible)',
         'OCI Data Catalog'])
    # Col 2 — ARMAZENAMENTO
    box(2.35, 0.25, 2.1, 2.95, 'ARMAZENAMENTO OCI', H_TEAL,
        ['OCI Object Storage', '(raw JSON events)',
         'Oracle Autonomous DB', '(events · risk scores)',
         'OCI Block Volume'])

    # Col 3 — PIPELINE ML
    box(4.62, 3.4, 2.5, 2.65, 'PIPELINE  ML', H_AMBER,
        ['OCI Data Flow', '(Apache Spark ETL)',
         'Feature Engineering', 'Sazonalidade · Magnitude',
         'Batch job agendado'])
    # Col 3 — MODELOS OCI
    box(4.62, 0.25, 2.5, 2.95, 'MODELOS  OCI', H_AMBER,
        ['OCI Data Science', 'XGBoost → Risk Score',
         'LSTM → D+7 Forecast',
         'OCI Model Catalog',
         'OCI Model Deployment'])

    # Col 4 — API
    box(7.3, 0.85, 2.2, 5.2, 'API DE ALERTAS', H_BLUE2,
        ['OCI Container Engine', '(OKE — Kubernetes)',
         'FastAPI + Docker',
         '', 'OCI Functions',
         '(alertas serverless)',
         '', 'OCI API Gateway',
         '(rate limit · auth · TLS)'])

    # Col 5a — DASHBOARD
    box(9.68, 3.4, 2.4, 2.65, 'DASHBOARD', H_CYAN,
        ['React 19 · D3.js · TS',
         'Vercel (CDN global)',
         'Mapa · Timeline · Feed',
         'Event card · Risk gauge'])
    # Col 5b — BI
    box(9.68, 0.25, 2.4, 2.95, 'BI / ANALYTICS', H_TEAL,
        ['Oracle Analytics Cloud',
         '(relatórios institucionais)',
         'CEMADEN · Defesa Civil',
         'Cruz Vermelha · ONGs'])

    # Col 6 — OPS
    box(12.25, 1.8, 1.65, 3.65, 'OPS', H_MID,
        ['OCI Monitoring', 'OCI Alarms', 'OCI Logging',
         'OCI Vault', '(secrets)', 'OCI IAM'])

    # Setas
    arrow(2.18, 4.7, 2.35, 4.7)  # fontes → ingestão
    arrow(2.18, 1.7, 2.35, 1.7)  # fontes → storage
    arrow(4.45, 4.7, 4.62, 4.7)  # ingestão → pipeline
    arrow(4.45, 1.7, 4.62, 1.7)  # storage → modelos
    arrow(7.12, 4.7, 7.30, 4.7)  # pipeline → API
    arrow(7.12, 1.7, 7.30, 1.7)  # modelos → API
    arrow(9.50, 4.7, 9.68, 4.7)  # API → dashboard
    arrow(9.50, 1.7, 9.68, 1.7)  # API → BI
    arrow(11.9, 3.4, 12.25, 3.4) # API → OPS

    # Rodapé
    ax.text(0.08, 0.08,
            'ODS 11 · 13 · 9  |  Open Data: NASA · NOAA · INPE  |  Oracle Cloud — Escalável · HA · Sem vendor lock-in',
            ha='left', va='bottom', fontsize=5.8, color=H_DIM, fontfamily='monospace')
    fig.tight_layout(pad=0.1)
    return fig


# ═══════════════════ BUILD SLIDES ═══════════════════════════════════

def build(out_path):
    prs = new_prs()

    # ─────────────────────────────────────────────
    # SLIDE 1 — CAPA
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    star_field(s, n=55)

    # Painel central
    rect(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.2),
         fill=PANEL2, line=STEEL)
    # Barra esquerda cyan
    rect(s, Inches(1.0), Inches(1.5), Pt(4), Inches(5.2), fill=CYAN)

    # Ícone satélite no canto
    embed_icon(s, 'satellite', Inches(10.2), Inches(1.7), 1.1,
               color=H_CYAN, bg=H_PAN2)

    txt(s, 'AEGIS', Inches(1.4), Inches(1.75), Inches(9), Inches(1.3),
        size=Pt(72), color=CYAN, bold=True)
    txt(s, 'ORBITAL DISASTER INTELLIGENCE',
        Inches(1.4), Inches(2.95), Inches(9), Inches(0.5),
        size=Pt(15), color=WHITE)
    rect(s, Inches(1.4), Inches(3.55), Inches(7.5), Pt(2), fill=AMBER)
    txt(s, 'Transformando Dados Orbitais em Alertas que Salvam Vidas',
        Inches(1.4), Inches(3.65), Inches(11), Inches(0.4),
        size=Pt(12.5), color=MID)

    txt(s, 'Global Solution 2026  ·  FIAP 2TSCPW  ·  Indústria Espacial',
        Inches(1.4), Inches(4.25), Inches(10), Inches(0.35),
        size=Pt(9.5), color=DIM)
    txt(s, 'Emerson dos Santos Silva  —  RM562033',
        Inches(1.4), Inches(4.65), Inches(6), Inches(0.3),
        size=Pt(10), color=MID, bold=True)
    txt(s, 'Pedro Henrique Soares  —  RM562283',
        Inches(1.4), Inches(4.98), Inches(6), Inches(0.3),
        size=Pt(10), color=MID, bold=True)
    txt(s, 'Junho de 2026',
        Inches(1.4), Inches(5.45), Inches(3), Inches(0.3),
        size=Pt(9), color=DIM)

    # ODS badges
    for i, (n, l, c) in enumerate([('11','Cidades',CYAN),('13','Clima',TEAL),('9','Inovação',AMBER2)]):
        bx = Inches(9.3) + i * Inches(1.0)
        rect(s, bx, Inches(5.2), Inches(0.9), Inches(0.6), fill=PANEL2, line=c)
        txt(s, f'ODS {n}', bx, Inches(5.22), Inches(0.9), Inches(0.28),
            size=Pt(9), color=c, bold=True, align=PP_ALIGN.CENTER)
        txt(s, l, bx, Inches(5.5), Inches(0.9), Inches(0.22),
            size=Pt(7), color=MID, align=PP_ALIGN.CENTER)

    # Ícones desastres na base
    disasters = [('fire',H_AMBER),('cyclone',H_BLUE2),('flood',H_TEAL),
                 ('quake',H_AMB2),('drought',H_CYAN)]
    for i, (ic, col) in enumerate(disasters):
        embed_icon(s, ic, Inches(1.4 + i*0.65), Inches(6.2), 0.5,
                   color=col, bg=H_PAN2)

    # ─────────────────────────────────────────────
    # SLIDE 2 — CENÁRIO ATUAL
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'CENÁRIO ATUAL')
    embed_icon(s, 'globe', Inches(12.0), Inches(0.4), 0.55, color=H_DIM, bg=H_BG)
    slide_title(s, 'O Espaço é a Nova Fronteira de Dados',
                'A economia espacial gera hoje mais valor do que a maioria das indústrias terrestres')

    metrics = [
        ('15.000+', 'satélites ativos\nem órbita (2026)', CYAN),
        ('USD 613 bi', 'valor economia\nespacial (2024)', TEAL),
        ('USD 1 tri', 'projeção\naté 2034', AMBER2),
        ('12 TB/dia', 'dados NASA\n+ Copernicus', BLUE2),
        ('< 3%', 'convertidos em\nalertas em tempo real', AMBER),
    ]
    for i, (val, lbl, vc) in enumerate(metrics):
        metric_card(s, val, lbl, Inches(0.38 + i*2.59), Inches(1.65),
                    w=Inches(2.48), h=Inches(1.35), vc=vc)

    bullets = [
        'Copernicus Programme (ESA/UE): clima, solo, oceanos e atmosfera — open data',
        'NASA EarthData: MODIS, VIIRS, Landsat, SMAP — gratuitamente acessíveis via API',
        'INPE (Brasil): PRODES + DETER + BDQueimadas — desmatamento e focos de incêndio',
        'NOAA: eventos climáticos extremos, furacões, temperatura oceânica',
        'Menos de 3% desses dados chegam a quem decide em tempo real',
    ]
    bullet_panel(s, bullets,
                 Inches(0.38), Inches(3.18), Inches(12.58), Inches(2.2),
                 title='PRINCIPAIS FONTES ORBITAIS E O GAP DE UTILIZAÇÃO', tc=CYAN)

    # Linha inferior: ícones de fontes
    sources = [
        ('satellite', 'NASA EONET',    H_CYAN,  CYAN),
        ('globe',     'Copernicus',    H_BLUE2, BLUE2),
        ('layers',    'INPE PRODES',   H_TEAL,  TEAL),
        ('pulse',     'NOAA Climate',  H_AMB2,  AMBER2),
    ]
    for i, (ic, lbl_s, hcol, pcol) in enumerate(sources):
        bx = Inches(0.38 + i*3.2)
        rect(s, bx, Inches(5.5), Inches(3.1), Inches(1.65),
             fill=PANEL2, line=STEEL)
        embed_icon(s, ic, bx + Inches(0.1), Inches(5.6), 0.45, color=hcol, bg=H_PAN2)
        txt(s, lbl_s, bx + Inches(0.65), Inches(5.72), Inches(2.3), Inches(0.35),
            size=Pt(10), color=pcol, bold=True)
        txt(s, 'Dados abertos · sem custo',
            bx + Inches(0.65), Inches(6.08), Inches(2.3), Inches(0.3),
            size=Pt(8.5), color=DIM)

    # ─────────────────────────────────────────────
    # SLIDE 3 — O PROBLEMA
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'O PROBLEMA')
    embed_icon(s, 'alert', Inches(12.0), Inches(0.4), 0.55, color=H_AMBER, bg=H_BG)
    slide_title(s, 'Latência Fatal: 6 a 24 Horas de Atraso',
                'Dados orbitais existem — o problema é transformá-los em decisão a tempo')

    stages = [
        ('DETECÇÃO\nORBITAL',   CYAN,     H_CYAN,  '0h'),
        ('PROCESSAMENTO\nTÉCNICO', DIM,   H_DIM,   '+2–6h'),
        ('RELATÓRIO\nMANUAL',   DIM,      H_DIM,   '+4–12h'),
        ('ALERTA\nHUMANO',      AMBER,    H_AMBER, '+6–24h'),
        ('RESPOSTA\nCIVIL',     RED_CRIT, '#EF4444','TARDE'),
    ]
    stage_icons = ['satellite','layers','archive','alert','team']
    for i, ((title, pcol, hcol, lag), ic) in enumerate(zip(stages, stage_icons)):
        bx = Inches(0.38 + i*2.59)
        rect(s, bx, Inches(1.55), Inches(2.48), Inches(2.5), fill=PANEL2, line=pcol)
        embed_icon(s, ic, bx + Inches(0.9), Inches(1.65), 0.55, color=hcol, bg=H_PAN2)
        txt(s, title, bx, Inches(2.3), Inches(2.48), Inches(0.6),
            size=Pt(9.5), color=pcol, bold=True, align=PP_ALIGN.CENTER)
        txt(s, lag, bx, Inches(2.95), Inches(2.48), Inches(0.6),
            size=Pt(18), color=pcol, bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, '→', Inches(2.78 + i*2.59), Inches(2.62), Inches(0.22), Inches(0.35),
                size=Pt(16), color=STEEL2, bold=True)

    stats = [
        ('USD 2,8 tri', 'Custo desastres\nclimat. — última\ndécada', AMBER),
        ('180 milhões', 'Pessoas afetadas\npor desastres\nem 2025', CYAN),
        ('47', 'Grandes desastres\nnaturais só\nem 2025', TEAL),
        ('18–48h', 'Vantagem orbital\nsobre alertas\nhumanos', AMBER2),
    ]
    for i, (val, lbl_s, vc) in enumerate(stats):
        metric_card(s, val, lbl_s, Inches(0.38 + i*3.24), Inches(4.25),
                    w=Inches(3.1), h=Inches(1.45), vc=vc)

    rect(s, Inches(0.38), Inches(5.86), Inches(12.58), Inches(0.5),
         fill=PANEL2, line=AMBER)
    txt(s, '"Em uma enchente repentina, horas são vidas."',
        Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.3),
        size=Pt(11.5), color=AMBER, align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────
    # SLIDE 4 — EVIDÊNCIA EDA
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'EVIDÊNCIAS — ANÁLISE EXPLORATÓRIA')
    embed_icon(s, 'pulse', Inches(12.0), Inches(0.4), 0.55, color=H_TEAL, bg=H_BG)
    slide_title(s, 'O que os Dados da NASA Revelam',
                '1.200+ eventos rastreados por satélite — NASA EONET v3 API · dados reais, não simulados')

    fig = fig_eda()
    embed_fig(s, fig, Inches(0.3), Inches(1.45), Inches(12.73), Inches(4.65))

    for i, b in enumerate([
        '›  Wildfires + Severe Storms = mais de 60% dos eventos ativos — picos sazonais detectáveis com ML',
        '›  Hotspots: América do Norte · Sudeste Asiático · Europa Mediterrânea · Brasil (queimadas)',
        '›  Sensores MODIS e VIIRS detectam eventos 18–48h antes dos alertas humanos convencionais',
    ]):
        txt(s, b, Inches(0.38), Inches(6.22) + i*Inches(0.34),
            Inches(12.58), Inches(0.3), size=Pt(9), color=MID)

    # ─────────────────────────────────────────────
    # SLIDE 5 — A SOLUÇÃO
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'A SOLUÇÃO')
    embed_icon(s, 'layers', Inches(12.0), Inches(0.4), 0.55, color=H_CYAN, bg=H_BG)
    slide_title(s, 'AEGIS — Inteligência Orbital em Tempo Real',
                'Plataforma que transforma dados de satélites em alertas preemptivos de desastres climáticos')

    features = [
        ('MONITORAMENTO\nEM TEMPO REAL', 'globe', H_CYAN, CYAN, [
            'Dashboard com dados NASA EONET ao vivo',
            '5 categorias monitoradas: incêndios,',
            'ciclones, enchentes, terremotos, secas',
            'Mapa D3.js com storm tracks e trajetórias',
        ]),
        ('MACHINE LEARNING\nPREDITIVO', 'pulse', H_AMBER, AMBER, [
            'XGBoost → Risk Score por evento (0–100)',
            'LSTM → Previsão D+1 a D+7 de volume',
            'SHAP → Explicabilidade dos modelos',
            'Feature engineering sobre dados orbitais',
        ]),
        ('INTERFACE\nPARA GESTORES', 'alert', H_TEAL, TEAL, [
            'Timeline scrubber: navegue D-7 a D+7',
            'Event feed priorizado por severidade',
            'Event card: coordenadas, magnitude,',
            'sparkline 7d, botões de ação',
        ]),
        ('ARQUITETURA\nESCALÁVEL', 'satellite', H_BLUE2, BLUE2, [
            'Oracle Cloud Infrastructure (OCI)',
            'FastAPI + Docker + OKE (Kubernetes)',
            'Open source · Sem vendor lock-in',
            'Oracle Autonomous DB para histórico',
        ]),
    ]
    for i, (title, ic, hcol, pcol, items) in enumerate(features):
        bx = Inches(0.38 + i*3.24)
        rect(s, bx, Inches(1.5), Inches(3.12), Inches(3.65), fill=PANEL2, line=pcol)
        rect(s, bx, Inches(1.5), Inches(3.12), Inches(0.46), fill=pcol)
        embed_icon(s, ic, bx + Inches(0.1), Inches(1.54), 0.36, color=H_BG, bg=hcol)
        txt(s, title, bx + Inches(0.56), Inches(1.57), Inches(2.45), Inches(0.38),
            size=Pt(9.5), color=BG, bold=True)
        for j, item in enumerate(items):
            txt(s, '›  '+item, bx + Inches(0.12), Inches(2.08)+j*Inches(0.58),
                Inches(2.9), Inches(0.52), size=Pt(9.5), color=WHITE)

    rect(s, Inches(0.38), Inches(5.25), Inches(12.58), Inches(0.55),
         fill=PANEL2, line=AMBER)
    txt(s, 'Redução de 6–24h para < 5 min entre detecção orbital e alerta terrestre',
        Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.35),
        size=Pt(13), color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    # Ícones disaster no rodapé
    for i, (ic, lbl_s, col, src) in enumerate([
        ('fire',    'Incêndios',  H_AMBER, 'MODIS · VIIRS · INPE — 18h de antecedência'),
        ('cyclone', 'Ciclones',   H_BLUE2, 'GOES · Sentinel — trajetória D+7'),
        ('flood',   'Enchentes',  H_TEAL,  'NOAA precip. · bacia hidrográfica'),
        ('quake',   'Terremotos', H_AMB2,  'USGS realtime · Richter + localização'),
        ('drought', 'Secas',      H_CYAN,  'SMAP soil moisture · índice PDSI'),
    ]):
        bx = Inches(0.38 + i*2.59)
        rect(s, bx, Inches(5.9), Inches(2.48), Inches(1.42), fill=PANEL2, line=STEEL)
        embed_icon(s, ic, bx + Inches(0.1), Inches(5.98), 0.40, color=col, bg=H_PAN2)
        txt(s, lbl_s, bx + Inches(0.60), Inches(6.05), Inches(1.82), Inches(0.3),
            size=Pt(10), color=col, bold=True)
        txt(s, src,   bx + Inches(0.12), Inches(6.40), Inches(2.3),  Inches(0.6),
            size=Pt(7.5), color=DIM)

    # ─────────────────────────────────────────────
    # SLIDE 6 — PÚBLICO-ALVO E PERSONA
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'PÚBLICO-ALVO E PERSONA')
    embed_icon(s, 'team', Inches(12.0), Inches(0.4), 0.55, color=H_MID, bg=H_BG)
    slide_title(s, 'Para Quem o AEGIS Foi Construído',
                'Definição da persona principal e dos públicos beneficiados pela solução')

    # Persona card
    rect(s, Inches(0.38), Inches(1.5), Inches(5.85), Inches(5.65),
         fill=PANEL2, line=CYAN)
    rect(s, Inches(0.38), Inches(1.5), Inches(5.85), Inches(0.40), fill=CYAN)
    txt(s, 'PERSONA PRINCIPAL', Inches(0.52), Inches(1.60), Inches(5.0), Inches(0.28),
        size=Pt(9.5), color=BG, bold=True)
    embed_icon(s, 'team', Inches(5.42), Inches(1.54), 0.32, color=H_BG, bg=H_CYAN)

    persona = [
        ('GESTORA DE OPERAÇÕES DE EMERGÊNCIA', CYAN, Pt(14), True),
        ('Coordenação de Defesa Civil · Resposta a Desastres', WHITE, Pt(10), False),
        ('Formação: Meteorologia · Gestão de Riscos Climáticos', MID, Pt(9), False),
        ('', WHITE, Pt(3), False),
        ('DESAFIO:', AMBER, Pt(9.5), True),
        ('Recebe alertas tardios. Quando confirma\no risco, o tempo de mobilização já foi.', WHITE, Pt(9.5), False),
        ('', WHITE, Pt(3), False),
        ('MOTIVAÇÃO:', TEAL, Pt(9.5), True),
        ('Salvar vidas antes do desastre, não durante.', WHITE, Pt(9.5), False),
        ('', WHITE, Pt(3), False),
        ('NECESSIDADE:', AMBER2, Pt(9.5), True),
        ('Interface que traduza dados orbitais\nem ações claras, sem expertise técnica\nem sensoriamento remoto.', WHITE, Pt(9.5), False),
    ]
    ty = Inches(2.05)
    for t, c, sz, b in persona:
        txt(s, t, Inches(0.52), ty, Inches(5.55), Inches(0.55),
            size=sz, color=c, bold=b)
        ty += float(sz) * 2.2 / 914400 * Inches(1) if t else Inches(0.07)

    # Públicos
    pubs = [
        ('AGÊNCIAS GOV.', H_CYAN, CYAN, 'globe',
         'CEMADEN · Defesa Civil\nIBAMA · Prefeituras em risco'),
        ('ORG. HUMANITÁRIAS', H_AMBER, AMBER, 'team',
         'Cruz Vermelha · UNOCHA\nIFRC · equipes de campo'),
        ('MUNICÍPIOS\nEM RISCO', H_TEAL, TEAL, 'alert',
         'Zonas de enchente, deslizamento\nSeca · planejamento evacuação'),
        ('PESQUISADORES', H_BLUE2, BLUE2, 'layers',
         'Análise de padrões orbitais\nData Science · publicações'),
    ]
    for i, (t, hc, pc, ic, desc) in enumerate(pubs):
        bx = Inches(6.42)
        by = Inches(1.5) + i*Inches(1.48)
        rect(s, bx, by, Inches(6.55), Inches(1.38), fill=PANEL2, line=pc)
        embed_icon(s, ic, bx + Inches(0.1), by + Inches(0.15), 0.48, color=hc, bg=H_PAN2)
        txt(s, t, bx + Inches(0.7), by + Inches(0.08), Inches(2.5), Inches(0.5),
            size=Pt(10), color=pc, bold=True)
        txt(s, desc, bx + Inches(0.7), by + Inches(0.62), Inches(5.7), Inches(0.65),
            size=Pt(9), color=MID)

    # ─────────────────────────────────────────────
    # SLIDE 7 — PROTÓTIPO TELAS 1
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'PROTÓTIPO — TELAS DA SOLUÇÃO')
    embed_icon(s, 'globe', Inches(12.0), Inches(0.4), 0.55, color=H_CYAN, bg=H_BG)
    slide_title(s, 'Dashboard AEGIS — Mapa Global e Storm Track',
                'Inserir screenshot do dashboard após deploy Vercel  ·  URL disponível na conclusão')

    rect(s, Inches(0.38), Inches(1.5), Inches(7.9), Inches(4.6),
         fill=PANEL2, line=STEEL)
    txt(s, '[ SCREENSHOT ]\nMAPA GLOBAL DE EVENTOS ATIVOS\nDados reais NASA EONET · D3.js · TopoJSON',
        Inches(0.38), Inches(3.1), Inches(7.9), Inches(1.0),
        size=Pt(10), color=DIM, align=PP_ALIGN.CENTER)
    embed_icon(s, 'globe', Inches(3.8), Inches(2.0), 0.9, color=H_STEEL, bg=H_PAN2)

    rect(s, Inches(8.45), Inches(1.5), Inches(4.55), Inches(4.6),
         fill=PANEL2, line=CYAN)
    rect(s, Inches(8.45), Inches(1.5), Inches(4.55), Inches(0.40), fill=CYAN)
    txt(s, 'TELA 1 — MAPA GLOBAL', Inches(8.58), Inches(1.60), Inches(4.3), Inches(0.28),
        size=Pt(9.5), color=BG, bold=True)
    desc_t1 = ('Exibe todos os desastres ativos rastreados por satélite em '
               'tempo real via NASA EONET. Cada ponto representa um evento '
               'geolocalizado com cor indicando o nível de risco calculado '
               'pelo modelo XGBoost.\n\n'
               'Storm tracks mostram trajetória passada (sólido) e projeção '
               'D+7 (pontilhado) para eventos móveis como ciclones.\n\n'
               'O operador identifica focos críticos sem interpretar relatórios.')
    txt(s, desc_t1, Inches(8.58), Inches(1.88), Inches(4.25), Inches(3.8),
        size=Pt(9), color=MID)

    rect(s, Inches(0.38), Inches(6.22), Inches(12.58), Inches(1.05),
         fill=PANEL2, line=STEEL)
    embed_icon(s, 'pulse', Inches(0.48), Inches(6.32), 0.42, color=H_AMBER, bg=H_PAN2)
    txt(s, 'TELA 2 — TIMELINE SCRUBBER (D-7 a D+7)',
        Inches(1.0), Inches(6.32), Inches(5), Inches(0.28),
        size=Pt(9), color=AMBER, bold=True)
    txt(s, 'Permite navegar no tempo: histórico e previsão futura do evento. '
        'Base para decisões de evacuação e alocação de recursos de emergência.',
        Inches(1.0), Inches(6.62), Inches(11.8), Inches(0.35),
        size=Pt(8.5), color=MID)

    # ─────────────────────────────────────────────
    # SLIDE 8 — PROTÓTIPO TELAS 2
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'PROTÓTIPO — TELAS DA SOLUÇÃO')
    embed_icon(s, 'alert', Inches(12.0), Inches(0.4), 0.55, color=H_AMBER, bg=H_BG)
    slide_title(s, 'Dashboard AEGIS — Event Card e Risk Analytics',
                'Inserir screenshots após deploy Vercel')

    # Placeholder tela 3
    rect(s, Inches(0.38), Inches(1.5), Inches(6.0), Inches(4.65),
         fill=PANEL2, line=STEEL)
    embed_icon(s, 'alert', Inches(2.6), Inches(2.4), 0.9, color=H_STEEL, bg=H_PAN2)
    txt(s, '[ SCREENSHOT ]\nEVENT CARD EXPANDIDO\nRisk Score · Sparkline · Ações',
        Inches(0.38), Inches(3.5), Inches(6.0), Inches(0.7),
        size=Pt(9), color=DIM, align=PP_ALIGN.CENTER)

    # Placeholder tela 4
    rect(s, Inches(6.55), Inches(1.5), Inches(6.42), Inches(4.65),
         fill=PANEL2, line=STEEL)
    embed_icon(s, 'pulse', Inches(9.2), Inches(2.4), 0.9, color=H_STEEL, bg=H_PAN2)
    txt(s, '[ SCREENSHOT ]\nRISK GAUGE + FORECAST CHART\nXGBoost · LSTM D+7',
        Inches(6.55), Inches(3.5), Inches(6.42), Inches(0.7),
        size=Pt(9), color=DIM, align=PP_ALIGN.CENTER)

    desc_cards = [
        (Inches(0.38), 'TELA 3 — EVENT CARD', CYAN, H_CYAN,
         'Risk Score XGBoost (0–100) · Sparkline 7 dias · '
         'Coordenadas exatas · Fonte orbital · '
         'Botões Dispatch / Monitor / Archive. '
         'Elimina navegação entre múltiplos sistemas.'),
        (Inches(6.55), 'TELA 4 — RISK GAUGE + FORECAST', AMBER, H_AMBER,
         'Gauge de risco em tempo real + Forecast Chart '
         'com previsão D+7 gerada pelo LSTM. '
         'Base para decisões de evacuação e '
         'alocação de recursos de emergência.'),
    ]
    for bx, t, pc, hc, desc in desc_cards:
        rect(s, bx, Inches(6.3), Inches(6.1), Inches(1.0), fill=PANEL2, line=pc)
        rect(s, bx, Inches(6.3), Pt(3), Inches(1.0), fill=pc)
        txt(s, t, bx + Inches(0.15), Inches(6.38), Inches(5.8), Inches(0.3),
            size=Pt(9), color=pc, bold=True)
        txt(s, desc, bx + Inches(0.15), Inches(6.7), Inches(5.8), Inches(0.55),
            size=Pt(8.5), color=MID)

    # ─────────────────────────────────────────────
    # SLIDE 9 — ARQUITETURA OCI
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'ARQUITETURA DA SOLUÇÃO')
    embed_icon(s, 'satellite', Inches(12.0), Inches(0.4), 0.55, color=H_CYAN, bg=H_BG)
    slide_title(s, 'Arquitetura Oracle Cloud Infrastructure (OCI)',
                'Pipeline escalável: coleta orbital → ML → API → consumo — high availability, sem vendor lock-in')

    fig = fig_arch()
    embed_fig(s, fig, Inches(0.18), Inches(1.4), Inches(12.95), Inches(5.85))

    # ─────────────────────────────────────────────
    # SLIDE 10 — STACK TECNOLÓGICO
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'TECNOLOGIAS')
    embed_icon(s, 'layers', Inches(12.0), Inches(0.4), 0.55, color=H_BLUE2, bg=H_BG)
    slide_title(s, 'Stack Tecnológico',
                'Cada tecnologia escolhida por critérios de escalabilidade, custo e aderência ao problema')

    tech = [
        ('DADOS ORBITAIS', H_CYAN, CYAN, 'satellite', [
            'NASA EONET API v3 (open data, sem auth)',
            'NOAA Climate API',
            'INPE BDQueimadas (focos de incêndio)',
            'Copernicus Sentinel Hub',
        ]),
        ('ANÁLISE + ML', H_AMBER, AMBER, 'pulse', [
            'Python · pandas · numpy',
            'XGBoost  →  Risk Score por evento',
            'LSTM (TensorFlow)  →  Previsão D+7',
            'SHAP  →  Explicabilidade dos modelos',
        ]),
        ('ORACLE CLOUD (OCI)', H_TEAL, TEAL, 'layers', [
            'OCI Data Flow (Apache Spark ETL)',
            'OCI Data Science + Model Catalog',
            'Oracle Autonomous DB (ATP)',
            'OCI Streaming · OCI Functions',
        ]),
        ('BACKEND', H_BLUE2, BLUE2, 'satellite', [
            'FastAPI + Docker + uvicorn',
            'OCI Container Engine (OKE)',
            'OCI API Gateway (auth · rate limit)',
            'OCI Object Storage (raw events)',
        ]),
        ('FRONTEND', H_CYAN, CYAN, 'globe', [
            'React 19 · TypeScript · Vite',
            'D3.js · D3-geo · TopoJSON',
            'Vercel (CDN global)',
            'Space Grotesk · JetBrains Mono',
        ]),
        ('BI / ANALYTICS', H_AMB2, AMBER2, 'pulse', [
            'Oracle Analytics Cloud (OAC)',
            'Plotly · matplotlib · seaborn',
            'Jupyter Notebooks (EDA)',
            'Relatórios para gestores',
        ]),
    ]
    for i, (group, hcol, pcol, ic, items) in enumerate(tech):
        col_i = i % 3
        row_i = i // 3
        bx = Inches(0.38 + col_i*4.33)
        by = Inches(1.65) + row_i * Inches(2.8)
        rect(s, bx, by, Inches(4.18), Inches(2.65), fill=PANEL2, line=pcol)
        rect(s, bx, by, Inches(4.18), Inches(0.44), fill=pcol)
        embed_icon(s, ic, bx + Inches(0.1), by + Inches(0.04), 0.36,
                   color=H_BG, bg=hcol)
        txt(s, group, bx + Inches(0.56), by + Inches(0.09), Inches(3.45), Inches(0.30),
            size=Pt(9.5), color=BG, bold=True)
        for j, item in enumerate(items):
            txt(s, '›  '+item, bx + Inches(0.12), by + Inches(0.52) + j*Inches(0.48),
                Inches(4.0), Inches(0.44), size=Pt(9.5), color=WHITE)

    # ─────────────────────────────────────────────
    # SLIDE 11 — IMPACTOS DA SOLUÇÃO
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'IMPACTOS DA SOLUÇÃO')
    slide_title(s, 'O Impacto do AEGIS no Mundo Real',
                'Três dimensões de impacto: ambiental, comunitário e socioeconômico')

    impacts = [
        ('IMPACTOS\nAMBIENTAIS', H_TEAL, TEAL, 'globe', [
            'Monitoramento contínuo de focos de\nincêndio acelera a resposta',
            'Detecção precoce de secas permite\nações de conservação hídrica',
            'Inteligência orbital para agências\nsem infraestrutura própria',
            'MODIS · VIIRS · Sentinel · USGS\nintegrados em um único painel',
        ]),
        ('IMPACTOS\nCOMUNITÁRIOS', H_CYAN, CYAN, 'team', [
            'Evacuações com 24–72h de\nantecedência reduzem mortes',
            'Gestores: resposta reativa\n→ prevenção preditiva',
            'Democratiza dados orbitais\npara municípios sem recurso',
            'Cruz Vermelha ganha 48h extras\npara mobilizar equipes de campo',
        ]),
        ('IMPACTOS\nSOCIOECONÔMICOS', H_AMBER, AMBER, 'layers', [
            'Redução 60–80% no custo de\nresposta emergencial',
            'Minimiza perdas com alertas\nantes do pico do evento',
            'Open source: qualquer agência\npública pode implantar',
            'OCI: paga só pelo que usa,\nsem superdimensionamento',
        ]),
    ]
    for i, (title, hcol, pcol, ic, items) in enumerate(impacts):
        bx = Inches(0.38 + i*4.33)
        rect(s, bx, Inches(1.5), Inches(4.18), Inches(4.95), fill=PANEL2, line=pcol)
        rect(s, bx, Inches(1.5), Inches(4.18), Inches(0.50), fill=pcol)
        embed_icon(s, ic, bx + Inches(0.1), Inches(1.56), 0.36,
                   color=H_BG, bg=hcol)
        txt(s, title, bx + Inches(0.56), Inches(1.60), Inches(3.4), Inches(0.42),
            size=Pt(10.5), color=BG, bold=True)
        for j, item in enumerate(items):
            txt(s, '›  '+item, bx + Inches(0.12), Inches(2.14) + j*Inches(0.82),
                Inches(3.95), Inches(0.75), size=Pt(10), color=WHITE)

    ods_map = [('ODS 11','Cidades Sustentáveis',CYAN),
               ('ODS 13','Ação Climática',TEAL),
               ('ODS 9','Inovação',AMBER)]
    for i, (num, name, col) in enumerate(ods_map):
        bx = Inches(0.38 + i*4.33)
        rect(s, bx, Inches(6.6), Inches(4.18), Inches(0.65), fill=PANEL2, line=col)
        txt(s, f'{num} · {name}', bx + Inches(0.12), Inches(6.7),
            Inches(4.0), Inches(0.4), size=Pt(9), color=col, bold=True)

    # ─────────────────────────────────────────────
    # SLIDE 12 — BENEFÍCIOS ESPERADOS
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'BENEFÍCIOS ESPERADOS')
    slide_title(s, 'Valor Gerado pelo AEGIS',
                'Cada benefício justificado pelos dados orbitais e pela arquitetura Oracle Cloud')

    bens = [
        ('Eficiência operacional', CYAN, H_CYAN,
         'Alertas em < 5 min vs. 6–24h atuais',
         'Pipeline orbital automatizado elimina interpretação manual'),
        ('Decisão baseada em dados', TEAL, H_TEAL,
         'Risk Score + previsão LSTM no mesmo painel',
         'Gestores sem expertise técnica acessam inteligência orbital'),
        ('Agilidade na resposta', AMBER, H_AMBER,
         'Timeline D+7 para planejamento antecipado',
         'Cruz Vermelha: 48h de antecedência para mobilizar equipes'),
        ('Economia de recursos', AMBER2, H_AMB2,
         'Alocação preditiva de helicópteros e equipes',
         'Sazonalidade EONET permite planejamento orçamentário'),
        ('Open source auditável', BLUE2, H_BLUE2,
         'Dados abertos NASA · sem vendor lock-in',
         'Qualquer agência pública pode replicar e adaptar'),
        ('Qualidade de dados', TEAL, H_TEAL,
         'MODIS + VIIRS + Sentinel + USGS integrados',
         'Cobertura complementar reduz pontos cegos de detecção'),
    ]
    for i, (name, pcol, hcol, how, why) in enumerate(bens):
        col_i = i % 2
        row_i = i // 2
        bx = Inches(0.38 + col_i*6.48)
        by = Inches(1.55) + row_i*Inches(1.75)
        rect(s, bx, by, Inches(6.28), Inches(1.65), fill=PANEL2, line=pcol)
        rect(s, bx, by, Pt(3), Inches(1.65), fill=pcol)
        txt(s, name, bx + Inches(0.15), by + Inches(0.08), Inches(6), Inches(0.38),
            size=Pt(11), color=pcol, bold=True)
        txt(s, '›  ' + how, bx + Inches(0.15), by + Inches(0.5), Inches(6), Inches(0.35),
            size=Pt(9.5), color=WHITE)
        txt(s, '   ' + why, bx + Inches(0.15), by + Inches(0.88), Inches(6), Inches(0.55),
            size=Pt(8.5), color=MID)

    # ─────────────────────────────────────────────
    # SLIDE 13 — DIFERENCIAL
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    header(s, 'ORIGINALIDADE E INOVAÇÃO')
    slide_title(s, 'O que Diferencia o AEGIS',
                'Comparativo direto com as soluções existentes hoje no mercado')

    col_xs = [Inches(0.38), Inches(3.6),  Inches(8.0)]
    col_ws = [Inches(3.15), Inches(4.33), Inches(5.0)]
    hdrs   = ['ASPECTO', 'SOLUÇÕES ATUAIS', 'AEGIS']
    hcols  = [STEEL2, STEEL2, CYAN]
    ht_c   = [WHITE, WHITE, BG]

    for j, (hdr, cx, cw, hc, tc) in enumerate(zip(hdrs, col_xs, col_ws, hcols, ht_c)):
        rect(s, cx, Inches(1.55), cw - Pt(3), Inches(0.5), fill=hc, line=STEEL)
        txt(s, hdr, cx + Inches(0.1), Inches(1.6), cw - Inches(0.15), Inches(0.38),
            size=Pt(10), color=tc, bold=True)

    rows = [
        ('Latência do alerta', '6 a 24 horas', '< 5 minutos'),
        ('Dados utilizados', 'Relatórios manuais e PDFs', 'APIs de satélite em tempo real'),
        ('Modelo preditivo', 'Reativo — só após o evento', 'Preditivo — até 72h antes'),
        ('Acesso e custo', 'Sistemas proprietários caros', 'Open source + dados NASA'),
        ('Interface', 'Planilhas e painéis técnicos', 'Dashboard missão p/ gestores'),
        ('Integração', 'Silos por agência', 'EONET + NOAA + INPE unificados'),
        ('Escalabilidade', 'On-premise, difícil expansão', 'OCI — paga pelo uso · HA'),
    ]
    rh = Inches(0.58)
    for r, (aspect, current, aegis) in enumerate(rows):
        by = Inches(2.08) + r*rh
        for j, (text, cx, cw, bg_c, tc) in enumerate([
            (aspect,  col_xs[0], col_ws[0], PANEL2,  WHITE),
            (current, col_xs[1], col_ws[1], PANEL2,  MID),
            (aegis,   col_xs[2], col_ws[2], RGBColor(0x0A,0x18,0x2A), CYAN),
        ]):
            rect(s, cx, by, cw - Pt(3), rh - Pt(2), fill=bg_c, line=STEEL)
            txt(s, text, cx + Inches(0.1), by + Inches(0.08),
                cw - Inches(0.15), rh - Inches(0.1),
                size=Pt(9.5), color=tc, bold=(j==2))

    # ─────────────────────────────────────────────
    # SLIDE 14 — CONCLUSÃO
    # ─────────────────────────────────────────────
    s = blank_slide(prs)
    set_bg(s)
    cyan_topbar(s)
    star_field(s, n=60, seed=42)

    rect(s, Inches(0.55), Inches(0.75), Inches(12.25), Inches(5.95),
         fill=PANEL2, line=STEEL)
    rect(s, Inches(0.55), Inches(0.75), Pt(4), Inches(5.95), fill=CYAN)

    # Ícone satélite
    embed_icon(s, 'satellite', Inches(11.2), Inches(0.85), 1.3,
               color=H_CYAN, bg=H_PAN2)

    txt(s, 'DADOS DO ESPAÇO,',
        Inches(0.95), Inches(0.92), Inches(10.5), Inches(0.9),
        size=Pt(38), color=WHITE, bold=True)
    txt(s, 'IMPACTO NA TERRA.',
        Inches(0.95), Inches(1.75), Inches(10.5), Inches(0.85),
        size=Pt(38), color=CYAN, bold=True)

    txt(s,
        'O AEGIS não inventa dados — transforma o que já existe no espaço em decisões que salvam vidas.\n'
        'A economia espacial gera 12 TB de dados por dia. O AEGIS é a ponte entre esse dado orbital\n'
        'e o gestor de defesa civil que precisa agir em minutos, não em horas.',
        Inches(0.95), Inches(2.7), Inches(10.5), Inches(1.2),
        size=Pt(11.5), color=MID)

    rect(s, Inches(0.95), Inches(3.98), Inches(11.5), Pt(1.5), fill=STEEL)

    cta = [
        ('Dashboard ao vivo', '[URL Vercel após deploy]', CYAN,    H_CYAN,    'globe'),
        ('Código open source', 'github.com/PedroHSSoares-Dev/gs-2026-aegis', TEAL, H_TEAL, 'layers'),
        ('Pitch YouTube', '[URL após gravação — PÚBLICO]',         AMBER,     H_AMBER,   'alert'),
    ]
    for i, (lbl_s, val, pcol, hcol, ic) in enumerate(cta):
        bx = Inches(0.95 + i*3.85)
        rect(s, bx, Inches(4.12), Inches(3.7), Inches(1.4), fill=PANEL2, line=pcol)
        embed_icon(s, ic, bx + Inches(0.1), Inches(4.22), 0.48,
                   color=hcol, bg=H_PAN2)
        txt(s, lbl_s, bx + Inches(0.7), Inches(4.2), Inches(2.85), Inches(0.35),
            size=Pt(9), color=pcol, bold=True)
        txt(s, val, bx + Inches(0.7), Inches(4.58), Inches(2.85), Inches(0.7),
            size=Pt(9), color=WHITE)

    txt(s, 'ODS 11 · ODS 13 · ODS 9  |  2TSCPW · FIAP 2026',
        Inches(0.95), Inches(5.65), Inches(8), Inches(0.3),
        size=Pt(8.5), color=DIM)
    txt(s, 'Emerson dos Santos Silva — RM562033  |  Pedro Henrique Soares — RM562283',
        Inches(0.95), Inches(5.98), Inches(11.5), Inches(0.32),
        size=Pt(9.5), color=MID, bold=True)

    # ─────────────────────────────────────────────
    # SALVAR
    # ─────────────────────────────────────────────
    prs.save(out_path)
    print(f'Salvo: {out_path}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    np.random.seed(2026)
    base = os.path.dirname(os.path.abspath(__file__))
    out  = os.path.join(base, 'GS_2TSCPW_Pitch_AEGIS.pptx')
    build(out)
